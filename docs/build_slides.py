# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=1.0"]
# ///
"""Ludoforge 프로젝트 소개 슬라이드를 python-pptx로 생성한다.

Marp(--pptx)는 각 슬라이드를 통째로 PNG로 박아 넣어 텍스트가 편집 불가하고
파일이 비대해진다. 이 스크립트는 네이티브 텍스트박스/도형/표를 직접 생성해
편집 가능하고 선명한 개발자 톤(다크) 슬라이드를 만든다.

빌드:
    uv run docs/build_slides.py            # docs/intro-slides.pptx 생성
미리보기(선택, LibreOffice 필요):
    soffice --headless --convert-to pdf --outdir docs docs/intro-slides.pptx

콘텐츠는 이 스크립트에 내장한다(SSOT). 내용을 바꾸려면 아래 SLIDES를 고친다.
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── 디자인 토큰 (GitHub Dark 계열 개발자 톤) ──────────────────────────────
BG = RGBColor(0x0D, 0x11, 0x17)  # 슬라이드 배경
PANEL = RGBColor(0x16, 0x1B, 0x22)  # 코드/표 패널
BORDER = RGBColor(0x30, 0x36, 0x3D)  # 구분선·테두리
TEXT = RGBColor(0xE6, 0xED, 0xF3)  # 본문
MUTED = RGBColor(0x8B, 0x94, 0x9E)  # 보조 텍스트·주석
BLUE = RGBColor(0x58, 0xA6, 0xFF)  # 주 강조
GREEN = RGBColor(0x3F, 0xB9, 0x50)  # 성공/긍정
GREEN_LT = RGBColor(0x7E, 0xE7, 0x87)  # 코드 내 강조
RED = RGBColor(0xFF, 0x7B, 0x72)  # 모순/실패
YELLOW = RGBColor(0xD2, 0x99, 0x22)  # 주의·라벨
PURPLE = RGBColor(0xBC, 0x8C, 0xFF)  # 보조 강조

# Windows 뷰잉 기준 안전 폰트(코드 내 한글은 PowerPoint가 자동 폴백)
FONT = "맑은 고딕"
FONT_MONO = "Consolas"

# 슬라이드 기하 (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.85)
CONTENT_W = SLIDE_W - MARGIN_X * 2
BODY_TOP = Inches(1.7)  # 제목 아래 본문 시작 y

# ── 인라인 마크업 파서: **굵게**, `코드` ─────────────────────────────────
_TOKEN = re.compile(r"\*\*(.+?)\*\*|`(.+?)`")


def _segments(text: str):
    """문자열을 (텍스트, 스타일) 조각으로 쪼갠다. 스타일: plain/bold/code."""
    out, last = [], 0
    for m in _TOKEN.finditer(text):
        if m.start() > last:
            out.append((text[last : m.start()], "plain"))
        if m.group(1) is not None:
            out.append((m.group(1), "bold"))
        else:
            out.append((m.group(2), "code"))
        last = m.end()
    if last < len(text):
        out.append((text[last:], "plain"))
    return out or [("", "plain")]


def _style_run(run, style: str, size: int, base_color: RGBColor):
    f = run.font
    f.size = Pt(size)
    if style == "bold":
        f.name = FONT
        f.bold = True
        f.color.rgb = TEXT
    elif style == "code":
        f.name = FONT_MONO
        f.color.rgb = BLUE
    else:
        f.name = FONT
        f.color.rgb = base_color


def _set_para_runs(para, text: str, size: int, color: RGBColor):
    for seg_text, style in _segments(text):
        run = para.add_run()
        run.text = seg_text
        _style_run(run, style, size, color)


# ── 저수준 헬퍼 ─────────────────────────────────────────────────────────
def _solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_autosize(tf):
    tf.word_wrap = True
    # auto_size = None: 박스 크기 고정(자동 리사이즈 끔)
    tf.auto_size = None


def _add_box(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    _no_autosize(tf)
    return box, tf


def _space_after(para, pt: float):
    para.space_after = Pt(pt)
    para.space_before = Pt(0)
    para.line_spacing = 1.15


# ── 슬라이드 골격 ────────────────────────────────────────────────────────
def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide


def add_footer(slide, page: int, total: int):
    # 좌하단 브랜드
    box, tf = _add_box(slide, MARGIN_X, Inches(7.02), Inches(4), Inches(0.35))
    p = tf.paragraphs[0]
    _set_para_runs(p, "Ludoforge", 10, MUTED)
    # 우하단 페이지
    box2, tf2 = _add_box(slide, SLIDE_W - Inches(2.0), Inches(7.02), Inches(1.15), Inches(0.35))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _set_para_runs(p2, f"{page} / {total}", 10, MUTED)


def add_title(slide, title: str, kicker: str | None = None):
    # 좌측 강조 바
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(0.72), Inches(0.10), Inches(0.62)
    )
    _solid(bar, BLUE)
    # 제목
    box, tf = _add_box(
        slide, MARGIN_X + Inches(0.28), Inches(0.6), CONTENT_W - Inches(0.28), Inches(0.9)
    )
    p = tf.paragraphs[0]
    if kicker:
        kp = tf.paragraphs[0]
        _set_para_runs(kp, kicker, 12, BLUE)
        _space_after(kp, 2)
        p = tf.add_paragraph()
    _set_para_runs(p, title, 30, TEXT)
    for r in p.runs:
        r.font.bold = True
    # 하단 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(1.48), CONTENT_W, Pt(1.2))
    _solid(line, BORDER)


# ── 블록 렌더러 (세로 커서 흐름) ─────────────────────────────────────────
def render_bullets(slide, items, y, width=None):
    """items: (text, level) 리스트. 반환: 다음 y."""
    width = width or CONTENT_W
    box, tf = _add_box(slide, MARGIN_X, y, width, Inches(5.0))
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if level == 0:
            marker = p.add_run()
            marker.text = "▸  "
            marker.font.name = FONT
            marker.font.size = Pt(18)
            marker.font.color.rgb = BLUE
            _set_para_runs(p, text, 18, TEXT)
            _space_after(p, 11)
        elif level == 1:
            p.level = 1
            marker = p.add_run()
            marker.text = "·  "
            marker.font.name = FONT
            marker.font.size = Pt(16)
            marker.font.color.rgb = MUTED
            _set_para_runs(p, text, 16, MUTED)
            _space_after(p, 7)
        else:  # level 2: 들여쓴 일반 문장(마커 없음)
            p.level = 1
            _set_para_runs(p, text, 16, TEXT)
            _space_after(p, 7)
    return y  # 본문은 박스가 충분히 큼


def render_note(slide, text, y, color=BLUE, width=None):
    """좌측 강조 막대가 붙은 콜아웃 한 줄."""
    width = width or CONTENT_W
    h = Inches(0.62)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, y, Inches(0.07), h)
    _solid(bar, color)
    box, tf = _add_box(slide, MARGIN_X + Inches(0.25), y, width - Inches(0.25), h)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _set_para_runs(p, text, 17, TEXT)
    return y + h + Inches(0.15)


def _color_code_line(p, line: str, lang: str):
    """간단한 코드 색칠: 주석/프롬프트/리포트 기호 강조."""
    if lang == "text":
        if line.startswith("❌"):
            _set_para_runs(p, line, 15, RED)
            return
        if line.lstrip().startswith("→"):
            _set_para_runs(p, line, 15, MUTED)
            return
        m = re.match(r"(\s*\[\d+\])(.*)", line)
        if m:
            r1 = p.add_run()
            r1.text = m.group(1)
            r1.font.name, r1.font.size, r1.font.color.rgb = FONT_MONO, Pt(15), YELLOW
            r2 = p.add_run()
            r2.text = m.group(2)
            r2.font.name, r2.font.size, r2.font.color.rgb = FONT_MONO, Pt(15), TEXT
            return
        _set_para_runs_mono(p, line, TEXT)
        return
    # yaml / bash: '#' 이후는 주석
    if lang == "bash" and line.lstrip().startswith("$"):
        idx = line.index("$")
        r0 = p.add_run()
        r0.text = line[:idx] + "$ "
        r0.font.name, r0.font.size, r0.font.color.rgb = FONT_MONO, Pt(15), GREEN
        rest = line[idx + 1 :].lstrip()
        r1 = p.add_run()
        r1.text = rest
        r1.font.name, r1.font.size, r1.font.color.rgb = FONT_MONO, Pt(15), TEXT
        return
    if lang == "lf" and "//" in line:  # 자체 문법(.lf) 주석은 `//`
        i = line.index("//")
        code, comment = line[:i], line[i:]
        if code:
            r0 = p.add_run()
            r0.text = code
            r0.font.name, r0.font.size, r0.font.color.rgb = FONT_MONO, Pt(15), TEXT
        r1 = p.add_run()
        r1.text = comment
        r1.font.name, r1.font.size, r1.font.color.rgb = FONT_MONO, Pt(15), GREEN_LT
        return
    if "#" in line:
        i = line.index("#")
        code, comment = line[:i], line[i:]
        if code:
            r0 = p.add_run()
            r0.text = code
            r0.font.name, r0.font.size, r0.font.color.rgb = FONT_MONO, Pt(15), TEXT
        r1 = p.add_run()
        r1.text = comment
        r1.font.name, r1.font.size, r1.font.color.rgb = FONT_MONO, Pt(15), GREEN_LT
        return
    _set_para_runs_mono(p, line, TEXT)


def _set_para_runs_mono(p, text: str, color: RGBColor):
    r = p.add_run()
    r.text = text if text else " "
    r.font.name, r.font.size, r.font.color.rgb = FONT_MONO, Pt(15), color


def render_code(slide, lines, lang, y, label=None, width=None):
    """코드 패널(둥근 사각형 + 모노스페이스). 반환: 다음 y."""
    width = width or CONTENT_W
    pad = Inches(0.22)
    line_h = Inches(0.275)
    body_h = line_h * len(lines)
    panel_h = body_h + pad * 2
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_X, y, width, panel_h)
    panel.adjustments[0] = 0.04
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = BORDER
    panel.line.width = Pt(1)
    panel.shadow.inherit = False

    box, tf = _add_box(slide, MARGIN_X + pad, y + pad, width - pad * 2, body_h)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        _color_code_line(p, ln, lang)

    next_y = y + panel_h + Inches(0.18)
    if label:
        lb, ltf = _add_box(slide, MARGIN_X, next_y, width, Inches(0.4))
        lp = ltf.paragraphs[0]
        _set_para_runs(lp, label, 14, MUTED)
        next_y += Inches(0.45)
    return next_y


def render_table(slide, head, rows, y, width=None):
    width = width or CONTENT_W
    cols = len(head)
    n = len(rows) + 1
    row_h = Inches(0.62)
    height = row_h * n
    gtbl = slide.shapes.add_table(n, cols, MARGIN_X, y, width, height)
    tbl = gtbl.table
    # python-pptx 기본 표 스타일 제거 → 직접 색칠
    tbl.first_row = False
    tbl.horz_banding = False
    widths = [int(width * 0.34), int(width * 0.66)] if cols == 2 else [int(width / cols)] * cols
    for c, w in enumerate(widths):
        tbl.columns[c].width = Emu(w)
    for r in range(n):
        tbl.rows[r].height = row_h

    def fill_cell(cell, text, color, *, bold=False, header=False, accent=None):
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL if header else BG
        cell.margin_left = Inches(0.18)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _set_para_runs(p, text, 16, accent or color)
        for run in p.runs:
            if bold or header:
                run.font.bold = True

    for c, htext in enumerate(head):
        fill_cell(tbl.cell(0, c), htext, BLUE, header=True)
    for r, row in enumerate(rows, start=1):
        accent = None
        if "Ludoforge" in row[0]:
            accent = GREEN_LT
        for c, val in enumerate(row):
            fill_cell(tbl.cell(r, c), val, TEXT, bold=(c == 0), accent=accent if c == 0 else None)
    return y + height + Inches(0.2)


# ── 콘텐츠 (SSOT) ────────────────────────────────────────────────────────
def build_title_slide(prs):
    slide = new_slide(prs)
    # 중앙 정렬 타이틀 블록
    box, tf = _add_box(slide, MARGIN_X, Inches(2.1), CONTENT_W, Inches(3.2))
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Ludoforge"
    r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(64), True, TEXT
    _space_after(p, 6)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "게임 기획 검증 툴킷"
    r2.font.name, r2.font.size, r2.font.color.rgb = FONT, Pt(26), BLUE
    _space_after(p2, 22)

    p3 = tf.add_paragraph()
    _set_para_runs(p3, "기획 룰을 **논리는 증명**, **정량은 추정**으로 검증하는 도구", 20, TEXT)
    _space_after(p3, 4)

    # 강조 막대
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(2.05), Inches(0.12), Inches(1.95)
    )
    _solid(bar, BLUE)
    # 좌측 막대만큼 타이틀 박스 들여쓰기
    box.left = MARGIN_X + Inches(0.32)

    # 하단 대상
    fb, ftf = _add_box(slide, MARGIN_X + Inches(0.32), Inches(5.9), CONTENT_W, Inches(0.5))
    fp = ftf.paragraphs[0]
    _set_para_runs(fp, "대상: 게임 기획자 · 프로그래머", 16, MUTED)


SLIDES = [
    {
        "title": "왜 이런 툴이 필요한가",
        "blocks": [
            (
                "bullets",
                [
                    ("기획자가 여럿이면, 각자 **합리적으로** 쓴 룰이 함께 두면 모순된다:", 0),
                    ('기획자 A — "전사 최대 HP = 레벨 × 100"', 1),
                    ('기획자 B — "모든 캐릭터 HP는 5000을 넘지 않는다"', 1),
                    ('기획자 C — "레벨 상한은 100"', 1),
                ],
            ),
            (
                "note",
                "레벨 51 전사는 HP 5100이어야 하는데 상한 5000과 충돌 → "
                "존재할 수 없는 상태가 조용히 생긴다",
                RED,
            ),
            (
                "bullets",
                [
                    ("룰이 수백 개면 사람 눈으로 모든 조합을 검토하는 건 **불가능**하다.", 0),
                ],
            ),
        ],
    },
    {
        "title": "기존 방식 vs Ludoforge",
        "blocks": [
            (
                "table",
                ["방식", "모순 발견"],
                [
                    ["사람 리뷰", "놓치기 쉬움 (조합 폭발)"],
                    ["시뮬레이션", "우연히 그 상태를 만나야 발견"],
                    ["Ludoforge", "모순의 존재 자체를 증명 (반례 없이도)"],
                ],
            ),
            (
                "note",
                "핵심 원칙: 판정은 사람·LLM이 아니라 **Z3(SMT solver)** 가 한다 — "
                '결정론이라 "거짓 일관성"을 환각하지 않는다',
                BLUE,
            ),
        ],
    },
    {
        "title": "DSL과 Z3",
        "kicker": "알아둘 개념 (1)",
        "blocks": [
            (
                "bullets",
                [
                    (
                        "**DSL** — 게임 룰을 산문 대신 **기계가 검증할 수 있는 자체 문법(`.lf`)**으로 작성.",
                        0,
                    ),
                    ("형식화하는 행위 자체가 숨은 가정을 드러낸다.", 2),
                    (
                        "**Z3 (SMT solver)** — `x + y <= 10 ∧ x > 3` 같은 산술 논리식을 푸는 도구.",
                        0,
                    ),
                    ("답은 셋 중 하나:", 2),
                    ("**sat** — 모든 제약을 만족하는 값이 존재 (예시 값을 줌)", 1),
                    ("**unsat** — 어떤 값으로도 전부 만족 불가 = **모순**", 1),
                    ("**unknown** — 시간초과/이론적 한계로 판단 못 함 (숨기지 않고 따로 보고)", 1),
                ],
            ),
        ],
    },
    {
        "title": "범인 룰과 도달성",
        "kicker": "알아둘 개념 (2)",
        "blocks": [
            (
                "bullets",
                [
                    (
                        "**unsat core (범인 룰)** — unsat일 때 Z3가 "
                        "**모순을 일으킨 최소 룰 집합**을 돌려준다.",
                        0,
                    ),
                    ('"이 룰들이 서로 싸운다"를 정확히 짚는다.', 2),
                    ("**핵심 통찰** — 룰을 그냥 다 모아 unsat을 물으면 **모순을 놓친다.**", 0),
                    ("Z3가 role=mage처럼 모순을 피해 가는 값을 골라버리기 때문.", 2),
                    (
                        '**명시적 도달성 검사** — "합법이라 여기는 상태(레벨 100 전사)를 '
                        '룰들이 봉쇄하는가?"',
                        0,
                    ),
                    ("봉쇄되면 모순.", 2),
                ],
            ),
        ],
    },
    {
        "title": "설치 — 준비물은 uv 하나",
        "blocks": [
            (
                "code",
                "bash",
                None,
                [
                    "# 1) uv 설치 (Python은 uv가 자동으로 받아온다)",
                    "#    https://docs.astral.sh/uv/getting-started/installation/",
                    "",
                    "# 2) Ludoforge 설치",
                    "uv tool install git+https://github.com/haje01/ludoforge.git",
                    "",
                    "# 3) 어디서나 사용",
                    "ludoforge check <룰 폴더>",
                ],
            ),
            ("note", "Windows 비개발자도 위 3단계면 끝 — 별도 Python 설치 불필요", GREEN),
        ],
    },
    {
        "title": "DSL 구조 — domain + constraints",
        "blocks": [
            (
                "code",
                "lf",
                None,
                [
                    "domain {                          // ① 변수와 그 범위 선언",
                    "    level: int 1..100",
                    "    hp:    int 0..",
                    "    role:  enum { warrior, mage, archer }",
                    "}",
                    "",
                    "constraint warrior_hp_formula:    // ② 지켜야 할 제약",
                    "    when role == warrior          // 조건 → Implies(when, then)",
                    "    then hp == level * 100",
                    "constraint global_hp_cap:",
                    "    then hp <= 5000",
                ],
            ),
            ("note", "각 룰은 `id`로 추적되어, 모순 시 어떤 룰이 범인인지 짚을 수 있다", BLUE),
        ],
    },
    {
        "title": "검증 사례 — 모순을 짚어낸다",
        "blocks": [
            (
                "code",
                "bash",
                None,
                [
                    "$ ludoforge check warrior.lf",
                ],
            ),
            (
                "code",
                "text",
                None,
                [
                    "❌ 모순 1건이 발견되었습니다.",
                    "",
                    "[1] role=warrior일 때 'level'은(는) 최대 50까지만",
                    "    도달 가능합니다 (선언 max=100).",
                    "    → 범인 룰: global_hp_cap, warrior_hp_formula",
                ],
            ),
            (
                "note",
                "어떤 조건에서 / 무엇이 봉쇄됐고 / 누가 범인인지 한국어로 보고 · "
                "모순이면 종료코드 1 → CI에서 PR마다 자동 차단",
                GREEN,
            ),
        ],
    },
    {
        "title": "팀 협업 — domain과 constraints 분리",
        "blocks": [
            (
                "bullets",
                [
                    ("공유 도메인 1개 + 기획자별 제약 파일. 디렉토리째 검사하면 병합된다.", 0),
                ],
            ),
            (
                "code",
                "text",
                None,
                [
                    "rules/",
                    "  _domain.lf     // 공유: 변수 정의만",
                    "  planner_a.lf   // 기획자 A: 전사 HP 공식 (constraints만)",
                    "  planner_b.lf   // 기획자 B: HP 상한 (constraints만)",
                ],
            ),
            (
                "note",
                "각 파일은 정상이어도, 합쳤을 때 생기는 **파일 간 모순**까지 잡아낸다 "
                "(import 불필요 · 변수 충돌은 오류로 보고)",
                BLUE,
            ),
        ],
    },
    {
        "title": "표현력 확장 — 2차",
        "kicker": "더 많은 룰을 형식화",
        "blocks": [
            (
                "bullets",
                [
                    (
                        "**불리언 상태 · 상호 배제** — `not (stealthed and attacking)`. "
                        "봉쇄된 상태를 찾는다. (D6)",
                        0,
                    ),
                    (
                        '**확률 · 실수(LRA)** — `type: real`로 "확률 합 = 1"을 직접. '
                        "`1/3`은 정확한 유리수. (D7)",
                        0,
                    ),
                    (
                        "**enum 타입 안전(EnumSort)** — 다른 enum이 같은 값 이름을 써도 안전. (D8)",
                        0,
                    ),
                    (
                        '**실수 끝점 도달성** — "[0,1]로 선언했는데 룰이 0.5로 막음" 류 '
                        "봉쇄를 잡는다. (D9)",
                        0,
                    ),
                    (
                        '**명시적 도달성 단언 `expect:`** — "두 스탯 동시 최대"처럼 '
                        "조합 도달성을 직접 검증. (D10)",
                        0,
                    ),
                ],
            ),
            (
                "note",
                "도달성 검사가 정수 → 불리언 · 실수 · 변수 조합으로 확장됐다 — "
                "핵심 원리(Z3 · unsat core)는 그대로",
                BLUE,
            ),
        ],
    },
    {
        "title": "예제 — enum · 불리언 · expect 함께",
        "kicker": "한 룰셋에서 같이 쓰기",
        "blocks": [
            (
                "code",
                "lf",
                None,
                [
                    "domain {",
                    "    role:      enum { rogue, mage }   // enum",
                    "    stealthed: bool                   // 불리언 상태",
                    "    attacking: bool",
                    "}",
                    "constraint stealth_mutex:",
                    "    then not (stealthed and attacking)     // 상호 배제",
                    "expect rogue_ambush:                       // 도달성 단언",
                    "    role == rogue and stealthed and attacking",
                ],
            ),
            (
                "code",
                "text",
                None,
                [
                    "❌ 모순 1건이 발견되었습니다.",
                    "[1] 기대 'rogue_ambush'가 충족되지 않습니다(도달 불가).",
                    "    → 범인 룰: stealth_mutex",
                ],
            ),
        ],
    },
    {
        "title": "동역학과 다중 백엔드 — 3차",
        "kicker": "시간과 확률을 더하다",
        "blocks": [
            (
                "bullets",
                [
                    ("정적 스냅샷을 넘어 **턴·이동·누적이 있는 동역학**을 검사한다.", 0),
                    ("`init` / `transitions`(가드·확률 분기·`x = expr` 대입) / `checks`를 DSL에 추가.", 2),
                    (
                        "**모델은 하나, 백엔드는 둘** — 시간과 확률은 다른 수학이라 "
                        "엔진을 나눈다. (D11)",
                        0,
                    ),
                    (
                        "**논리 백엔드(Z3/BMC, `ludoforge bmc`)** — 도달성·불변식·데드락을 "
                        "k 스텝 증명, 반례 경로 제시.",
                        1,
                    ),
                    (
                        "**정량 추정(Monte Carlo, `ludoforge sim`)** — 직업별 승률·분포를 "
                        "표집 추정(신뢰구간). PRISM은 소형 모델 교차검증 테스트 오라클. (D19·D23)",
                        1,
                    ),
                ],
            ),
            (
                "note",
                "보드게임 *던전!*(WotC)을 논리·정량 양쪽으로 검증한 게 동기 — "
                "같은 `.lf` 한 벌에서.",
                BLUE,
            ),
        ],
    },
    {
        "title": "정량 추정 — Monte Carlo (sim) — 4차",
        "kicker": "증명에서 추정으로 (D19)",
        "blocks": [
            (
                "bullets",
                [
                    ("PRISM은 **상태 폭발**이 천장 — 정확한 확률이 필요한 곳이 곧 못 푸는 곳.", 0),
                    ("정량 경로를 **표집 추정(Monte Carlo, `ludoforge sim`)** 으로 옮긴다.", 0),
                    ("승률·기대 길이·분포를 **신뢰구간**과 함께 추정 (증명 아님).", 1),
                    ("real·고차원·큰 범위도 상태폭발 없이 — 미관측 사건은 rule-of-three 상한.", 1),
                    ("**PRISM은 교차검증 오라클** — 소형 모델에서 증명기가 추정기를 검정.", 0),
                ],
            ),
            (
                "note",
                "존재·건전성은 *증명*(Z3/BMC), 정량 크기는 *추정*(sim) — "
                "한 DSL에서 각자의 수학으로.",
                GREEN,
            ),
        ],
    },
    {
        "title": "최근 진화 — 자체 DSL · 정책 · 단순화",
        "kicker": "5~7차",
        "blocks": [
            (
                "bullets",
                [
                    (
                        "**sim 정책(`pref`)** — 플레이어 *선택*에 확률을 줘 표집한다. 규칙과 전략을 "
                        '분리해 "이 정책에선 승률이?"를 추정("Pmax 아님" 라벨). (D20)',
                        0,
                    ),
                    (
                        "**자체 문법 `.lf`** — YAML을 벗고 전용 DSL로 승격. `=`(대입)·`==`(비교)를 "
                        "문법이 구분, 프라임 표기 제거. IR은 불변이라 백엔드 무회귀. (D21·D22)",
                        0,
                    ),
                    (
                        "**구조 단순화** — PRISM을 *테스트 전용 교차검증 오라클*로 내리고, 사용자 "
                        "백엔드를 `check`·`bmc`·`sim` 셋으로. 던전 예제도 하나로 통합. (D23)",
                        0,
                    ),
                ],
            ),
            (
                "note",
                '한 모델, 두 질문 — `bmc`는 "이길 길이 존재하는가"(증명), `sim`은 '
                '"이 정책에선 실제 얼마인가"(추정). `pref`는 sim만 본다.',
                GREEN,
            ),
        ],
    },
    {
        "title": "한계 (현재)",
        "blocks": [
            (
                "bullets",
                [
                    (
                        "**선형 산술(LIA·LRA) 중심** — `level*100`은 OK, "
                        "변수 × 변수(비선형)는 느리거나 판단 불가(우회 안내).",
                        0,
                    ),
                    (
                        "**비목표** — 밸런스 *재미* 평가, 런타임 서버 검증은 아니다(기획 단계 도구).",
                        0,
                    ),
                    (
                        "단, 밸런스 *튜닝*(승률·분포)은 4차에서 sim **추정**의 목표로 들어왔다. (D19)",
                        2,
                    ),
                    (
                        "**아직(후속)** — 실수 범위 정밀 분석(Optimize), "
                        "DTMC 정적 사전 검사, CI PR 코멘트 연동.",
                        0,
                    ),
                ],
            ),
        ],
    },
]


def build_content_slide(prs, spec):
    slide = new_slide(prs)
    add_title(slide, spec["title"], spec.get("kicker"))
    y = BODY_TOP
    for block in spec["blocks"]:
        kind = block[0]
        if kind == "bullets":
            y = render_bullets(slide, block[1], y)
            # bullets 박스는 높이를 자동으로 잡지 못하므로 다음 블록 위치를 항목 수로 추정
            est = sum(0.5 if lvl == 0 else 0.42 for _, lvl in block[1])
            y = y + Inches(est) + Inches(0.1)
        elif kind == "note":
            color = block[2] if len(block) > 2 else BLUE
            y = render_note(slide, block[1], y, color)
        elif kind == "code":
            _, lang, label, lines = block
            y = render_code(slide, lines, lang, y, label)
        elif kind == "table":
            y = render_table(slide, block[1], block[2], y)
    return slide


def build_closing_slide(prs):
    slide = new_slide(prs)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(2.5), Inches(0.12), Inches(2.0)
    )
    _solid(bar, BLUE)
    box, tf = _add_box(slide, MARGIN_X + Inches(0.32), Inches(2.5), CONTENT_W, Inches(2.4))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "감사합니다"
    r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(54), True, TEXT
    _space_after(p, 8)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Q & A"
    r2.font.name, r2.font.size, r2.font.color.rgb = FONT, Pt(28), BLUE
    _space_after(p2, 28)
    p3 = tf.add_paragraph()
    _set_para_runs(p3, "GitHub: `github.com/haje01/ludoforge`", 16, MUTED)
    _space_after(p3, 4)
    p4 = tf.add_paragraph()
    _set_para_runs(p4, "문서: `docs/concepts.md` · `examples/`", 16, MUTED)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    for spec in SLIDES:
        build_content_slide(prs, spec)
    build_closing_slide(prs)

    total = len(prs.slides._sldIdLst)
    for i, slide in enumerate(prs.slides, start=1):
        if i == 1:
            continue  # 표지에는 푸터 생략
        add_footer(slide, i, total)

    out = Path(__file__).parent / "intro-slides.pptx"
    prs.save(str(out))
    print(f"생성 완료: {out}  (슬라이드 {total}장)")


if __name__ == "__main__":
    main()
